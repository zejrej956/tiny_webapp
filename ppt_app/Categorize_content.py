from pptx import Presentation
from file_to_list import FileToList
import os

class Categorize_content:
	"""docstring for Categorize_content"""
	def __init__(self, content, filename, slide_template):

		self.filename = filename
		self.list_of_text = content.split("\n")

		
		print "been here..core"
		slide_design = slide_template

		self.prs = Presentation(slide_design)

		# create front slide of the presentation
		# Note slide and self.slide variables are not equivalent
		title_slide_layout = self.prs.slide_layouts[0]
		slide = self.prs.slides.add_slide(title_slide_layout)
		title = slide.shapes.title
		subtitle = slide.placeholders[1]
		title.text = self.list_of_text[0]
		subtitle.text = self.list_of_text[1]

	def convert_to_ppt(self):
		for text in self.list_of_text[2:]:
			tab_level = len(text) - len(text.strip("\t"))
			# print tab_level
			if tab_level==0:
				self.create_new_bullet_slide()
				self.add_title(text)
			else:	
				self.add_content(text, tab_level)

	def create_new_bullet_slide(self):
		bullet_slide_layout = self.prs.slide_layouts[1]
		self.slide = self.prs.slides.add_slide(bullet_slide_layout)

	def add_title(self, title):
		self.shapes = self.slide.shapes
		title_shape = self.shapes.title
		self.body_shape = self.shapes.placeholders[1]
		title_shape.text = title

		self.tf = self.body_shape.text_frame 

	def add_content(self, content, level):
		level -= 1
		p = self.tf.add_paragraph() 
		p.text = content.strip("\t")
		p.level = level

	def save_ppt(self):
		path_file = 'static/media/' + self.filename + '.pptx'
		self.prs.save(path_file)
